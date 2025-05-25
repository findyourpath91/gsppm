import os
import json
import pytz
import traceback
from datetime import datetime
from flask import Flask, request, jsonify
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from pptx import Presentation
from google.auth import default
from googleapiclient.http import MediaFileUpload

app = Flask(__name__)

# Google Drive service (authentication and setup)
from google.oauth2 import service_account

def authenticate_drive():
    """Authenticate and return the Google Drive service instance."""
    try:
        credentials_path = "/app/credentials/service-account.json"
        credentials = service_account.Credentials.from_service_account_file(
            credentials_path,
            scopes=["https://www.googleapis.com/auth/drive"]
        )
        return build('drive', 'v3', credentials=credentials)
    except Exception as e:
        raise RuntimeError(f"Authentication failed: {e}")

def get_folder_id_from_request(request_data):
    """Extract and validate folder ID."""
    folder_id = request_data.get('folderId')
    if not folder_id or not isinstance(folder_id, str):
        raise ValueError("Invalid or missing 'folderId' in request data.")
    return folder_id

def get_template_path_from_folder(drive_service, folder_id):
    """Locate .pptm file in specified Google Drive folder."""
    query = f"'{folder_id}' in parents and mimeType='application/vnd.ms-powerpoint.presentation.macroEnabled.12'"
    results = drive_service.files().list(q=query).execute()
    files = results.get('files', [])
    if not files:
        raise FileNotFoundError(f"No PowerPoint (.pptm) templates found in folder {folder_id}.")
    return files[0]['id']

def set_public_permissions(drive_service, file_id):
    """Set the permission of the file to 'Anyone with the link can view'."""
    try:
        drive_service.permissions().create(
            fileId=file_id,
            body={'role': 'reader', 'type': 'anyone'}
        ).execute()
        print(f"File {file_id} is now publicly accessible.")
    except HttpError as e:
        # Specific HttpError for set_public_permissions is handled within the function
        # The task asks to add this handling in process_request, so we'll do it there too.
        # For now, this internal print will remain.
        print(f"An error occurred while setting permissions: {e}")
        raise # Re-raise to be caught by process_request or the main handler

def upload_file_to_drive(drive_service, file_path, folder_id):
    """Upload a file to Google Drive."""
    file_metadata = {
        'name': os.path.basename(file_path),
        'parents': [folder_id]
    }
    media = MediaFileUpload(file_path, mimetype='application/vnd.ms-powerpoint.presentation.macroEnabled.12')
    try:
        uploaded_file = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        return uploaded_file['id']
    except Exception as e:
        # This will be caught by the specific HttpError or generic Exception in process_request
        raise RuntimeError(f"Failed to upload file: {e}")


def create_folder_in_drive(drive_service, folder_name, parent_folder_id=None):
    """Create a folder in Google Drive, ensuring no duplication."""
    query = f"name='{folder_name}' and mimeType='application/vnd.google-apps.folder'"
    if parent_folder_id:
        query += f" and '{parent_folder_id}' in parents"
    
    try:
        results = drive_service.files().list(q=query, fields="files(id, name)").execute()
        existing_folders = results.get('files', [])
        
        if existing_folders:
            return existing_folders[0]['id']
        
        folder_metadata = {
            'name': folder_name,
            'mimeType': 'application/vnd.google-apps.folder',
            'parents': [parent_folder_id] if parent_folder_id else []
        }
        folder = drive_service.files().create(body=folder_metadata, fields='id').execute()
        return folder['id']
    except Exception as e:
        raise RuntimeError(f"Failed to create or find folder '{folder_name}': {e}")

def load_data_from_request_json(request_data):
    """Load questions and answers from request JSON data."""
    data = request_data.get('data', [])
    questions_answers = []

    for entry in data:
        question = entry.get('question', '')
        answer = entry.get('answer', '')

        if isinstance(question, (int, float)):
            question = str(question)
        if isinstance(answer, (int, float)):
            answer = str(answer)

        questions_answers.append((question, answer))

    return questions_answers

def update_presentation(questions_answers, pptx_file, output_pptx):
    """Update placeholders in PowerPoint with provided question-answer pairs."""
    prs = Presentation(pptx_file)
    qa_index = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and qa_index < len(questions_answers):
                text = shape.text.lower()
                if "question" in text:
                    shape.text_frame.clear()
                    shape.text_frame.text = questions_answers[qa_index][0]
                elif "answer" in text:
                    shape.text_frame.clear()
                    shape.text_frame.text = questions_answers[qa_index][1]
                    qa_index += 1
    prs.save(output_pptx)

@app.route('/', methods=['GET'])
def health_check():
    """Provide a basic health check and welcome message."""
    return jsonify({"status": "Application is running", "message": "Welcome to the PPTX processing API"}), 200

@app.route('/process', methods=['POST'])
def process_request():
    """Handle incoming POST requests with PowerPoint data."""
    try:
        request_data = request.get_json()
        print(f"Received JSON payload: {request_data}")
        if not request_data:
            return jsonify({"error": "Invalid or missing JSON payload"}), 400

        folder_id = get_folder_id_from_request(request_data)
        questions_answers = load_data_from_request_json(request_data)
        if len(questions_answers) < 1:
            return jsonify({"error": "No question-answer data provided"}), 400

        cst = pytz.timezone("America/Chicago")
        timestamp = datetime.now(cst).strftime("%b %d %I %M %p")
        temp_dir = "/app/temp"
        os.makedirs(temp_dir, exist_ok=True)
        
        output_pptx_1 = os.path.join(temp_dir, f'Round 1_{timestamp}.pptm')
        output_pptx_2 = os.path.join(temp_dir, f'Round 2_{timestamp}.pptm')

        try:
            drive_service = authenticate_drive()
        except Exception as e:
            print(f"Error during Google Drive authentication: {e}")
            raise

        games_folder_id = create_folder_in_drive(drive_service, "Create Games", parent_folder_id=folder_id) # Assuming this can also raise, covered by main try-except
        
        try:
            template_file_id = get_template_path_from_folder(drive_service, folder_id)
        except FileNotFoundError as e:
            print(f"Error finding template file: {e}")
            raise
        except HttpError as e:
            print(f"Google Drive API error while getting template path: {e}")
            raise

        # Download the template file from Drive
        template_local_path = os.path.join(temp_dir, "template.pptm")
        try:
            drive_request = drive_service.files().get_media(fileId=template_file_id)
            with open(template_local_path, "wb") as f:
                f.write(drive_request.execute())
        except HttpError as e:
            print(f"Error downloading template file {template_file_id} from Drive: {e}")
            raise

        # Update presentations with the first and second sets of questions
        try:
            update_presentation(questions_answers[:25], template_local_path, output_pptx_1)
        except Exception as e:
            print(f"Error updating presentation {output_pptx_1}: {e}")
            raise
        try:
            update_presentation(questions_answers[25:], template_local_path, output_pptx_2)
        except Exception as e:
            print(f"Error updating presentation {output_pptx_2}: {e}")
            raise

        # Upload the updated presentations to Google Drive
        file_id_1 = None # Initialize in case of early error
        file_id_2 = None # Initialize in case of early error
        try:
            file_id_1 = upload_file_to_drive(drive_service, output_pptx_1, games_folder_id)
        except HttpError as e:
            print(f"Google Drive API error while uploading file {output_pptx_1}: {e}")
            raise
        except Exception as e:
            print(f"Error uploading file {output_pptx_1} to Drive: {e}")
            raise
        
        try:
            file_id_2 = upload_file_to_drive(drive_service, output_pptx_2, games_folder_id)
        except HttpError as e:
            print(f"Google Drive API error while uploading file {output_pptx_2}: {e}")
            raise
        except Exception as e:
            print(f"Error uploading file {output_pptx_2} to Drive: {e}")
            raise

        # Set public permissions for the uploaded files
        if file_id_1:
            try:
                set_public_permissions(drive_service, file_id_1)
            except HttpError as e:
                print(f"Google Drive API error while setting permissions for file {file_id_1}: {e}")
                raise
        if file_id_2:
            try:
                set_public_permissions(drive_service, file_id_2)
            except HttpError as e:
                print(f"Google Drive API error while setting permissions for file {file_id_2}: {e}")
                raise

        return jsonify({
            "success": True,
            "files": [
                {"file_name": os.path.basename(output_pptx_1), "file_id": file_id_1},
                {"file_name": os.path.basename(output_pptx_2), "file_id": file_id_2}
            ]
        }), 200

    except Exception as e:
        traceback.print_exc() # This was added in the previous step
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8111)
